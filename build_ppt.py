#!/usr/bin/env python3
"""
Ellectric 项目技术讲解 PPT 生成器
深色主题，专业配色
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════ 配色方案 ═══════════════════════
BG_DARK    = RGBColor(0x0D, 0x0D, 0x0D)    # 极深黑背景
BG_SECTION = RGBColor(0x1A, 0x1A, 0x2E)    # 深蓝黑分区背景
ACCENT     = RGBColor(0x00, 0xD4, 0xAA)     # 青绿色主强调
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)     # 紫色次强调
ACCENT3    = RGBColor(0x06, 0xB6, 0xD4)     # 蓝色
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
ORANGE     = RGBColor(0xF9, 0x73, 0x16)
RED        = RGBColor(0xEF, 0x44, 0x44)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout

def add_bg(slide, color=BG_DARK):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_line(slide, x1, y1, x2, y2, color=ACCENT, width=2):
    """添加直线"""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = Pt(width)

def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    """添加强调色条"""
    return add_shape_bg(slide, left, top, width, height, color)

def add_title_bar(slide, title, subtitle=None):
    """统一标题栏"""
    add_accent_bar(slide, Inches(0.8), Inches(0.4), Inches(0.06), Inches(0.5), ACCENT)
    add_text(slide, Inches(1.0), Inches(0.35), Inches(11), Inches(0.6),
             title, font_size=30, bold=True)
    if subtitle:
        add_text(slide, Inches(1.0), Inches(0.85), Inches(11), Inches(0.4),
                 subtitle, font_size=14, color=MID_GREY)
    # 底部分隔线
    add_shape_bg(slide, Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.01), DARK_GREY)

def add_card(slide, left, top, width, height, title, content_lines,
             title_color=ACCENT, border_color=None, bg_color=BG_SECTION):
    """添加卡片组件"""
    # 背景
    card = add_shape_bg(slide, left, top, width, height, bg_color)
    # 顶边强调线
    add_accent_bar(slide, left, top, width, Inches(0.04), border_color or title_color)
    # 标题
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.35),
             title, font_size=16, bold=True, color=title_color)
    # 内容
    y_offset = Inches(0.55)
    for line in content_lines:
        add_text(slide, left + Inches(0.2), top + y_offset, width - Inches(0.4), Inches(0.25),
                 line, font_size=11, color=LIGHT_GREY)
        y_offset += Inches(0.22)

def add_page_number(slide, num):
    """页码"""
    add_text(slide, Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3),
             str(num), font_size=10, color=MID_GREY, align=PP_ALIGN.RIGHT)

def add_bullet_list(slide, left, top, width, items, font_size=13, color=LIGHT_GREY):
    """带点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(6)
    return tf

# ═══════════════════════ SLIDE 1: 标题 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, BG_DARK)

# 顶部装饰线
add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

# 大标题
add_text(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
         "Ellectric", font_size=72, bold=True, color=WHITE)

# 副标题
add_text(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.6),
         "AI + 电力交易技术学习平台", font_size=28, color=ACCENT)

# 描述
add_text(slide, Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.8),
         "从公开数据到智能交易 — 端到端电力市场 AI 技术闭环",
         font_size=16, color=MID_GREY)

# 技术标签
tags = ["负荷预测", "电价预测", "电力市场仿真", "强化学习", "回测引擎", "SHAP 可解释"]
tag_y = Inches(4.8)
tag_x = Inches(1.5)
for tag in tags:
    tag_shape = add_shape_bg(slide, tag_x, tag_y, Inches(1.6), Inches(0.35), DARK_GREY)
    tag_shape.fill.fore_color.rgb = RGBColor(0x25, 0x25, 0x25)
    add_text(slide, tag_x, tag_y + Inches(0.05), Inches(1.6), Inches(0.3),
             tag, font_size=10, color=ACCENT, align=PP_ALIGN.CENTER)
    tag_x += Inches(1.8)

# 底部信息
add_text(slide, Inches(1.5), Inches(6.3), Inches(3), Inches(0.3),
         "Python 3.11+  |  XGBoost  |  Lasso  |  PPO/SAC/TD3  |  FastAPI  |  LangChain",
         font_size=10, color=MID_GREY)

# ═══════════════════════ SLIDE 2: 项目概述 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "项目概述", "什么是 Ellectric？目标是什么？")

add_bullet_list(slide, Inches(1.0), Inches(1.5), Inches(11),
    [
        "Ellectric = Electric + Lecture，一个动手实践型 AI + 电力交易学习平台",
        "目标：跑通「公开电力数据接入 → 负荷/电价预测 → 电力市场仿真 → 自动交易策略」端到端技术闭环",
        "非生产系统，教育/学习用途 — 每个模块都有详细中英双语注释和设计决策说明",
        "渐进式学习设计：从简单的持续法基线 → XGBoost → LEAR → 强化学习 → 完整 API 服务",
        "数据来源：Our World in Data (OWID) 全球能源数据 + 中国现货电价数据 (ZionLuo)",
    ],
    font_size=14,
)

# 三核心理念
cards = [
    ("🧪 动手实践", "10 个 Jupyter notebooks\n逐步骤学习每个环节"),
    ("📐 工程严谨", "TimeSeriesSplit + 防泄漏\nScaler fit-on-train-only"),
    ("🔬 可解释性", "SHAP + LEAR 系数\n模型不是黑盒"),
]
for i, (title, content) in enumerate(cards):
    left = Inches(1.0 + i * 3.9)
    add_card(slide, left, Inches(3.8), Inches(3.5), Inches(2.2), title, content.split("\n"),
             title_color=ACCENT2 if i == 1 else ACCENT)

add_page_number(slide, 2)

# ═══════════════════════ SLIDE 3: 四阶段演进 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "四阶段渐进式演进", "从数据到智能 — Phase 1 → Phase 4 技术路线图")

phases = [
    ("Phase 1", "数据管道 + XGBoost 预测", ACCENT,
     ["OWID 数据自动拉取（三级回退）", "数据清洗管道（缺失/异常/时区）", "三级渐进式特征工程 (Tier 1-3)",
      "XGBoost 负荷预测 (MAE/RMSE)", "TimeSeriesSplit 防时序泄露"]),
    ("Phase 2", "电价预测 + 市场仿真", ACCENT2,
     ["LEAR 电价预测 (Lasso L1 正则化)", "中国现货电价数据接入", "ASSUME 多智能体市场仿真",
      "有功出清 + 利润结算", "3 种预设场景(夏季高峰/大风)"]),
    ("Phase 3", "RL 交易 + 回测 + 可解释", ACCENT3,
     ["Gymnasium 电力交易环境", "PPO / TD3 / SAC 三种 RL 算法", "多策略回测引擎 (BacktestRunner)",
      "SHAP 特征重要性瀑布图", "Diebold-Mariano 统计检验"]),
    ("Phase 4", "API / CLI / LLM 接口", ORANGE,
     ["FastAPI REST 服务 (6 端点)", "Typer CLI 命令行工具", "LangChain + DeepSeek 对话 Agent",
      "SSE 流式 Chat 接口", "三明治架构：API/CLI/LLM → Service → Pipeline"])
]

for i, (title, subtitle, color, items) in enumerate(phases):
    left = Inches(0.5 + i * 3.15)
    add_card(slide, left, Inches(1.6), Inches(2.95), Inches(5.0),
             f"{title}: {subtitle}", items, title_color=color)

add_page_number(slide, 3)

# ═══════════════════════ SLIDE 4: 完整流程架构图 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "完整流程架构图", "端到端数据流 — 从原始数据到自动交易")

# 架构层
layers = [
    ("接入层", ACCENT, ["FastAPI REST API", "Typer CLI", "LangChain Agent"]),
    ("服务层", ACCENT2, ["handlers.py (4函数)", "schemas.py (Pydantic v2)", "延迟导入安全"]),
    ("交易层", ACCENT3, ["ElectricityMarketEnv", "PPO/SAC/TD3 Agent", "BacktestRunner"]),
    ("预测层", ORANGE, ["XGBoost → 负荷预测", "LEAR (Lasso) → 电价预测", "Persistence 基线对比"]),
    ("特征层", GREEN, ["Tier 1: 5个核心特征", "Tier 2: 节假日+周滞后", "Tier 3: 滚动统计+循环编码"]),
    ("数据层", RED, ["OWIDChinaLoader", "PriceDataLoader", "clean_data() 清洗管道"]),
]

for i, (name, color, items) in enumerate(layers):
    y = Inches(1.5 + i * 0.95)
    # 层标签
    add_shape_bg(slide, Inches(0.5), y, Inches(1.5), Inches(0.75), color)
    add_text(slide, Inches(0.5), y + Inches(0.2), Inches(1.5), Inches(0.4),
             name, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 箭头
    if i < len(layers) - 1:
        add_text(slide, Inches(1.15), y + Inches(0.65), Inches(0.3), Inches(0.3),
                 "▼", font_size=12, color=color, align=PP_ALIGN.CENTER)
    # 组件
    x = Inches(2.3)
    for item in items:
        card = add_shape_bg(slide, x, y + Inches(0.12), Inches(3.3), Inches(0.5), BG_SECTION)
        add_text(slide, x + Inches(0.15), y + Inches(0.22), Inches(3.0), Inches(0.35),
                 item, font_size=11, color=LIGHT_GREY)
        x += Inches(3.5)

# 底部数据流
add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
         "OWID GitHub → DataLoader → Cleaner → FeatureEngineer → XGBoost/LEAR → Gym Env → RL Agent → Backtest → SHAP → API/CLI/LLM",
         font_size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

add_page_number(slide, 4)

# ═══════════════════════ SLIDE 5: 数据接入层 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "数据接入层 — 电力数据从哪来？", "OWID 全球能源数据 + 中国现货电价数据")

# 左侧：OWID
add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
         "🌍 OWIDChinaLoader — 全球能源数据",
         [
             "数据源: Our World in Data (owid-energy-data.csv, ~25MB)",
             "三级回退策略: CDN → GitHub raw → 本地Parquet缓存",
             "流式解析: csv.DictReader 逐行过滤 iso_code='CHN'",
             "时间范围: 2000-2025, 频率: 年级",
             "TWh → 日均 MW: generation_twh × 1e6 / 365 / 24",
             "提取字段: 发电量/用电量/煤电/气电/风电/光伏/水电/核电",
             "",
             "核心列: timestamp (UTC) + load_mw (MW) + region",
             "抽象基类 DataLoader (ABC):",
             "  load_data(start, end) → DataFrame",
             "  get_metadata() → {source, version, rows}",
         ])

# 右侧：电价
add_card(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2),
         "⚡ PriceDataLoader — 中国现货电价",
         [
             "数据源: ZionLuo/Electricity-Price-Forecasting",
             "文件格式: .xlsx (Excel 文件)",
             "中文列名 → 英文列名自动映射 (7组映射)",
             "",
             "7 列输出:",
             "  timestamp  — 时间 (1h 粒度)",
             "  price_da   — 日前价格 (元/MWh)",
             "  price_rt   — 实时价格 (元/MWh)",
             "  load_mw    — 负荷 (MW)",
             "  wind_mw    — 风电出力 (MW)",
             "  solar_mw   — 光伏出力 (MW)",
             "  tie_line_mw — 省间联络线功率 (MW)",
             "",
             "独立加载器，不继承 DataLoader (列集不同)",
         ])

add_page_number(slide, 5)

# ═══════════════════════ SLIDE 6: 数据清洗 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "数据清洗管道 — 脏数据 → 干净数据", "clean_data() 四步清洗 + validate_schema() 质量检查")

steps = [
    ("1. 列验证", "检查 REQUIRED_COLUMNS\n= {timestamp, load_mw}",
     "缺少必须列直接报错\nFail Fast 原则"),
    ("2. 缺失值填充", "ffill() + bfill()\n时序连续性",
     "利用电力数据时间连续性\n前向填充优先，后向兜底"),
    ("3. 异常值检测", "IQR 方法 (Tukey's fences)\n1.5x 倍数",
     "⚠️ 只报告不删除\n尖峰=信号(极端天气/事件)"),
    ("4. 时区标准化", "无时区 → UTC\n非 UTC → 转 UTC",
     "所有时间戳统一 UTC\n数据合约强制要求"),
]

for i, (title, detail, reason) in enumerate(steps):
    left = Inches(0.5 + i * 3.15)
    # 步骤标题
    add_text(slide, left + Inches(0.2), Inches(1.6), Inches(2.8), Inches(0.3),
             title, font_size=18, bold=True, color=ACCENT)
    # 详情
    add_card(slide, left, Inches(2.0), Inches(2.95), Inches(2.0), detail, detail.split("\n"),
             title_color=ACCENT)
    # 原因
    add_card(slide, left, Inches(4.2), Inches(2.95), Inches(1.5), "设计原理", reason.split("\n"),
             title_color=ACCENT2, bg_color=RGBColor(0x20, 0x20, 0x30))

# 底部：质量评分
add_card(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0),
         "📊 get_data_quality_score() — 0-100 数据质量评分",
         ["缺失率 (40分) + 时序连续性 (30分) + 异常值比率 (30分)  →  综合质量评分 → 日志告警阈值 60分"],
         title_color=GREEN)

add_page_number(slide, 6)

# ═══════════════════════ SLIDE 7: 特征工程 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "特征工程 — 渐进式三级特征体系", "将领域知识编码为机器学习可理解的数字特征")

# 负荷预测特征 (左侧)
add_text(slide, Inches(0.5), Inches(1.5), Inches(5), Inches(0.3),
         "🔮 负荷预测特征 (FeatureEngineer)", font_size=16, bold=True, color=ACCENT)

tiers_load = [
    ("Tier 1 · 5 特征", GREEN, "hour, day_of_week, month, is_weekend, lag_24h",
     "核心: 日内周期 + 季节 + 持续法基础"),
    ("Tier 2 · 7 特征", ACCENT3, "+ is_holiday, lag_168h",
     "中级: 节假日模式 + 周周期"),
    ("Tier 3 · 11 特征", ACCENT2, "+ rolling_mean_24h, rolling_std_24h, hour_sin, hour_cos",
     "高级: 滚动统计 + sin/cos 循环编码 (保持 23→0 连续性)"),
]

for i, (name, color, cols, desc) in enumerate(tiers_load):
    y = Inches(2.0 + i * 1.2)
    add_shape_bg(slide, Inches(0.5), y, Inches(0.08), Inches(0.9), color)
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.35), name, font_size=14, bold=True, color=color)
    add_text(slide, Inches(0.8), y + Inches(0.35), Inches(5.5), Inches(0.3), cols, font_size=10, color=LIGHT_GREY)
    add_text(slide, Inches(0.8), y + Inches(0.6), Inches(5.5), Inches(0.25), desc, font_size=9, color=MID_GREY)

# 电价预测特征 (右侧)
add_text(slide, Inches(7.0), Inches(1.5), Inches(5), Inches(0.3),
         "💰 电价预测特征 (LEAR add_price_features)", font_size=16, bold=True, color=ACCENT2)

tiers_price = [
    ("Tier 1 · 6 特征", GREEN, "hour, day_of_week, month, is_weekend, lag_24h_price, lag_168h_price"),
    ("Tier 2 · 11 特征", ACCENT3, "+ lag_24h_load, lag_24h_wind, lag_24h_solar, rolling_mean/std_24h_price"),
    ("Tier 3 · 14 特征", ACCENT2, "+ hour_sin, hour_cos, price_trend_7d (168h 滚动均值)"),
]

for i, (name, color, cols) in enumerate(tiers_price):
    y = Inches(2.0 + i * 1.2)
    add_shape_bg(slide, Inches(7.0), y, Inches(0.08), Inches(0.9), color)
    add_text(slide, Inches(7.3), y, Inches(2.5), Inches(0.35), name, font_size=14, bold=True, color=color)
    add_text(slide, Inches(7.3), y + Inches(0.35), Inches(5.5), Inches(0.5), cols, font_size=10, color=LIGHT_GREY)

# 底部关键设计说明
add_card(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2),
         "🔑 关键设计",
         ["循环编码 (hour_sin/cos): 将 0-23 映射到单位圆，保持 23:00 和 00:00 在数值上的相邻关系",
          "级联补齐: 调用 Tier 3 时自动检测并补齐 Tier 1-2 缺失特征，调用者无需关心依赖"],
         title_color=ACCENT)
add_page_number(slide, 7)

# ═══════════════════════ SLIDE 8: XGBoost 负荷预测 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "负荷预测引擎 — XGBoost + TimeSeriesSplit", "为什么选 XGBoost？梯度提升原理 + 防时序泄露设计")

# 算法原理
add_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.0),
         "🧠 XGBoost 梯度提升原理",
         [
             "1. 从最简单的预测开始 (如均值)",
             "2. 计算残差 = 实际值 - 当前预测",
             "3. 训练一棵新决策树来预测这些残差",
             "4. 把新树的预测加到原来的预测上",
             "5. 重复 2-4，每次新树修正前一轮错误",
             "",
             "最终预测 = Tree₁ + Tree₂ + ... + Treeₙ",
             "",
             "参数: n_estimators=100, max_depth=6",
             "      learning_rate=0.1, subsample=0.8",
             "      colsample_bytree=0.8"],
         title_color=GREEN)

# 防泄漏设计
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.0),
         "🛡️ 防时序泄露 — 关键保证",
         [
             "TimeSeriesSplit(n_splits=5, gap=24):",
             "  训练: [2015...2017] → 测试: [2018]",
             "  训练: [2015...2018] → 测试: [2019]",
             "  永远过去预测未来!",
             "",
             "gap=24 参数: 训练/测试间隔 24 小时",
             "  防止 lag_24h 跨越训练/测试边界",
             "",
             "StandardScaler 只在训练集 fit:",
             "  fold 内 scaler.fit(X_train)",
             "  → transform(X_test)",
             "  绝不 fit_transform() 全量数据!"],
         title_color=RED)

# 对比：为什么不是深度学习
add_card(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.2),
         "✅ XGBoost 在电力负荷预测的优势",
         [
             "• 内置特征重要性 → 训练完就知道哪些特征有效",
             "• 原生处理缺失值 → 电力数据常有缺失",
             "• L1/L2 正则化防过拟合 → 比深度学习更稳健",
             "• CPU 优化 → 不需要 GPU，速度快",
             "• Kaggle 能源预测竞赛长期霸榜 → 工业验证"],
         title_color=GREEN)

# 持续法基线
add_card(slide, Inches(7.0), Inches(4.8), Inches(5.8), Inches(2.2),
         "📏 持续法基线 — 衡量一切的标尺",
         [
             "persistence_forecast(df):",
             "  forecast[t] = actual[t - 24h]",
             "",
             "为什么是 24 小时？",
             "  电力负荷有极强日周期 (Diurnal Cycle)",
             "  昨天下午 3 点 ≈ 今天下午 3 点",
             "",
             "如果 XGBoost MAE > 持续法 MAE",
             "  → 特征工程或模型选择有问题"],
         title_color=ORANGE)

add_page_number(slide, 8)

# ═══════════════════════ SLIDE 9: LEAR 电价预测 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "电价预测引擎 — LEAR (Lasso Estimated AutoRegressive)", "为什么用 Lasso 而不是 XGBoost？线性模型的可解释性优势")

# Lasso 原理
add_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.2),
         "📐 Lasso (L1 正则化) 原理",
         [
             "损失函数 = MSE + α × Σ|βᵢ|",
             "",
             "α (alpha) 控制正则化强度:",
             "  α=0   → 普通最小二乘, 可能过拟合",
             "  α 适中 → 部分系数压缩到零 (自动特征选择)",
             "  α 过大 → 所有系数为零 (欠拟合)",
             "",
             "LEAR = Lasso Estimated AutoRegressive",
             "来源: Lago et al., 2021, Applied Energy",
             "",
             "特征重要性 = 系数绝对值 |βᵢ|",
             "  正系数 (蓝色): 特征 ↑ → 价格 ↑",
             "  负系数 (红色): 特征 ↑ → 价格 ↓",
             "  未显示的特征: 系数被 L1 压缩为零"],
         title_color=ACCENT2)

# LEAR vs XGBoost
add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.2),
         "⚖️ LEAR vs XGBoost — 互补关系",
         [
             "XGBoost (负荷预测):",
             "  树模型 → 自动捕捉非线性交互",
             "  黑盒 → 难以解释每个特征的边际影响",
             "  适合: 负荷预测 (强非线性日周期)",
             "",
             "LEAR/Lasso (电价预测):",
             "  线性模型 → 系数 = 边际影响",
             "  白盒 → 每个特征的经济含义清晰",
             "  适合: 电价预测 (LEAR 精度经同行评审验证)",
             "",
             "在电价预测中 LEAR 通常 ≥ 复杂 DL 模型"],
         title_color=ACCENT3)

#
add_card(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0),
         "🔬 LEAR 特征系数可视化 — plot_coefficients()",
         [
             "lag_24h_price: 最强预测因子 — 昨天同时刻价格是最好的起点 (日周期锚定)",
             "lag_168h_price: 上周同时刻 — 捕捉周模式 (周末 vs 工作日电价差)",
             "lag_24h_load: 负荷滞后 — 负荷增加通常推高电价 (供需关系)",
             "hour_sin/cos: 小时循环编码 — 保持早晚高峰的连续过渡",
             "price_trend_7d: 7日价格趋势 — 捕捉中长期价格方向"],
         title_color=ACCENT)

add_page_number(slide, 9)

# ═══════════════════════ SLIDE 10: ASSUME 电力市场仿真 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "电力市场仿真 — ASSUME 框架 + 内置回退引擎", "多智能体竞价 → Merit Order 出清 → 结算")

# 仿真流程
add_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.5),
         "🏭 仿真流程",
         [
             "1. 读取 YAML 配置文件",
             "   default / summer_peak / wind_high",
             "",
             "2. 构建发电机组:",
             "   按燃料类型+边际成本分组",
             "   煤(50GW) 风(20GW) 光(15GW) 气(10GW) 储(5GW)",
             "",
             "3. 生成负荷曲线:",
             "   正弦日负荷 (峰15:00, 谷03:00)",
             "   叠加高斯噪声模拟真实波动",
             "",
             "4. Merit Order 经济调度:",
             "   新能源优先 (边际成本≈0)",
             "   → 火电按成本从低到高",
             "   → 出清价 = 最后一台机组边际成本"],
         title_color=ORANGE)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.5),
         "📊 输出与场景",
         [
             "4 个输出文件:",
             "  clearing_prices.csv   — 每小时出清价",
             "  dispatch.csv          — 各机组调度量",
             "  agent_profits.csv     — 各参与者利润",
             "  simulation_metadata.json — 仿真元信息",
             "",
             "3 种预设场景:",
             "  default:     煤50 + 风20 + 光15 + 气10 + 储5 (GW)",
             "  summer_peak: 需求 96GW, 增加气电",
             "  wind_high:   风电 30GW, 减少煤电",
             "",
             "ASSUME 框架 vs 内置引擎:",
             "  已安装 → assume.World 多智能体仿真",
             "  未安装 → run_builtin_simulation()"],
         title_color=GREEN)

# 价格接受者模型
add_card(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
         "💡 关键简化 — 价格接受者模型",
         [
             "本项目的仿真采用价格接受者 (Price Taker) 假设: 单个交易者投标量不影响出清价格",
             "出清价来自历史数据 (price_da) 或仿真计算 — 这简化了市场博弈，聚焦于智能体投标策略学习",
             "真实电力市场通常是寡头博弈 → 更复杂的均衡求解或强化学习多智能体博弈 (超出学习范围)"],
         title_color=ORANGE)

add_page_number(slide, 10)

# ═══════════════════════ SLIDE 11: 强化学习交易 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "强化学习交易环境 — ElectricityMarketEnv + RL Agent", "Gymnasium 强化学习 → PPO/TD3/SAC → 自动投标")

# 观测/动作空间
add_card(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(3.0),
         "👁️ 观测空间 (Dict)",
         [
             "load_forecast_24h   (24,)",
             "  → XGBoost 预测未来 24h 负荷",
             "price_forecast_24h  (24,)",
             "  → LEAR 预测未来 24h 电价",
             "time_features        (4,)",
             "  → hour+day_of_week sin/cos",
             "price_history_168h   (168,)",
             "  → 过去 7 天电价历史",
             "account_state        (2,)",
             "  → [现金, 进度%]"],
         title_color=ACCENT)

add_card(slide, Inches(4.8), Inches(1.5), Inches(4.0), Inches(3.0),
         "🎯 动作空间 + 出清逻辑",
         [
             "动作: Box(0, 1, (24,))",
             "  归一化投标量 ∈ [0, 1]",
             "  实际MW = action × max_capacity",
             "",
             "出清逻辑 (价格接受者):",
             "  cleared = min(bid_mw, actual_load)",
             "",
             "P&L = -|bid - actual| × price / 1000",
             "",
             "3 种内置奖励函数:",
             "  profit_only  — 仅总盈亏",
             "  risk_adjusted — 盈亏 - 1×σ",
             "  volume_penalty — 盈亏 - 过量惩罚"],
         title_color=ACCENT3)

add_card(slide, Inches(9.1), Inches(1.5), Inches(3.7), Inches(3.0),
         "🤖 三种 RL 算法",
         [
             "PPO (Proximal Policy Optimization):",
             "  On-policy, 稳定性好, 最推荐",
             "",
             "TD3 (Twin Delayed DDPG):",
             "  Off-policy, 连续动作空间优化",
             "  双 Q 网络减少过估计",
             "",
             "SAC (Soft Actor-Critic):",
             "  Off-policy, 最大熵探索",
             "  自动温度参数调节",
             "",
             "统一接口: BaseRLAgent (ABC)",
             "工厂: RLAgentFactory.create()"],
         title_color=ACCENT2)

# 训练流程
add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2),
         "🔄 训练流程 — RL Agent 学习投标策略",
         [
             "step 1: env.reset() → 初始观测              step 3: 出清 + 计算奖励              step 5: 终止？→ 下一回合",
             "step 2: agent.predict(obs) → 24 维投标动作   step 4: env.step(action) → 下一观测    重复 10000+ timesteps → 收敛"],
         title_color=ACCENT)

add_page_number(slide, 11)

# ═══════════════════════ SLIDE 12: 回测引擎 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "回测引擎 — 历史数据回放 + 多策略对比", "BacktestRunner: 在历史数据上「重演」交易策略")

# 三种基线策略
add_card(slide, Inches(0.5), Inches(1.5), Inches(4.0), Inches(2.5),
         "📏 三种基线策略",
         [
             "持续法 (Persistence):",
             "  bid[t] = 24h前实际负荷",
             "  代表传统调度员的直觉",
             "",
             "均值法 (Mean):",
             "  bid[t] = 过去168h (7天) 均值",
             "  代表简单统计方法",
             "",
             "先知 (Oracle):",
             "  bid[t] = 实际负荷 (完美预见)",
             "  代表理论上界 — 不可能超越"],
         title_color=GREEN)

add_card(slide, Inches(4.8), Inches(1.5), Inches(4.0), Inches(2.5),
         "🔄 回放流程",
         [
             "1. 确定策略 (基线 or RL Agent)",
             "2. 对齐负荷+价格数据时间范围",
             "3. 创建 ElectricityMarketEnv",
             "4. 逐 24h 块回放:",
             "   action = strategy(env, t)",
             "   → env.step(action)",
             "   → 记录每笔交易",
             "5. 输出 hourly records DataFrame"],
         title_color=ACCENT3)

add_card(slide, Inches(9.1), Inches(1.5), Inches(3.7), Inches(2.5),
         "📊 多策略对比指标",
         [
             "compare(results) 返回:",
             "  总收益 Total Return",
             "  夏普比率 Sharpe Ratio",
             "     = mean/σ × √8760",
             "  胜率 Win Rate",
             "     = P&L>0 的比例",
             "  最大回撤 Max Drawdown",
             "  交易次数"],
         title_color=ACCENT2)

# 可视化说明
add_card(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.7),
         "📈 累计 P&L 叠加图 — plot_comparison()",
         [
             "多策略累计盈亏曲线叠加在同一图表中，纵轴=累计P&L，横轴=时间",
             "Oracle (先知): 用真实负荷投标 → 完美 → 理论上界 (不可能被超越)",
             "Persistence (持续法): t-24h 作为今日投标 → 传统调度基准 → 智能体至少应该优于它",
             "RL Agent (PPO/TD3/SAC): 训练后自主投标 → 在 Oracle 和 Persistence 之间 → 越接近 Oracle 越好",
             "颜色: 蓝(先知) 橙(持续法) 绿(PPO) 红(TD3) 紫(SAC) — 多头对比一目了然"],
         title_color=ORANGE)

add_page_number(slide, 12)

# ═══════════════════════ SLIDE 13: SHAP + 统计检验 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "模型可解释性 + 统计检验", "SHAP 博弈论解释 + Diebold-Mariano 显著性检验")

add_card(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.5),
         "🎯 SHAP — Shapley Additive Explanations",
         [
             "基于博弈论 Shapley 值:",
             "  prediction = base_value + Σ shap_valueᵢ",
             "",
             "TreeExplainer (XGBoost):",
             "  精确计算每棵树中每个特征的边际贡献",
             "  输出: 水平柱状瀑布图",
             "  红色 → 特征推高预测",
             "  蓝色 → 特征拉低预测",
             "",
             "LinearExplainer (LEAR/Lasso):",
             "  利用回归系数直接计算 SHAP 值",
             "",
             "feature_importance_ranking():",
             "  跨模型 (XGBoost vs LEAR) 特征排名对比"],
         title_color=ACCENT)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(3.5),
         "📐 Diebold-Mariano + Giacomini-White 检验",
         [
             "DM 检验 (Diebold & Mariano, 1995):",
             "  比较两个预测模型的损失差异",
             "  H₀: 两模型预测精度无显著差异",
             "  p < 0.05 → 拒绝 H₀ → 显著差异",
             "",
             "GW 检验 (Giacomini & White, 2006):",
             "  DM 检验的推广 — 条件预测能力",
             "  同时检查无条件 + 条件预测准确性",
             "",
             "对比: 中国 LEAR vs 5 个基准",
             "  EPEX-BE/FR/DE, NordPool, PJM",
             "结论: 统计量化\"LEAR 是否真正领先\""],
         title_color=ACCENT2)

add_card(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.7),
         "💡 为什么需要 SHAP + 统计检验？",
         [
             "MAE 告诉你模型有多准，SHAP 告诉你模型为什么这样预测 — 两者互补",
             "SHAP waterfall 图可以解释单个时刻的预测: \"为什么预测 14:00 负荷这么高？→ 因为 lag_24h=高 + hour_sin=正\"",
             "DM/GW 检验提供统计置信度: \"LEAR 的 MAE 比 XGBoost 低 5%，这个差异是显著的还是随机波动？\""],
         title_color=ACCENT)

add_page_number(slide, 13)

# ═══════════════════════ SLIDE 14: API/CLI/LLM 三层接口 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "三层接口 — FastAPI + CLI + LangChain Agent", "三明治架构: API/CLI/LLM → Service → Pipeline")

# FastAPI
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(4.0),
         "🌐 FastAPI REST API",
         [
             "6 个端点:",
             "  GET  /health    → 健康检查",
             "  POST /predict   → 负荷/电价预测",
             "  POST /simulate  → 电力市场仿真",
             "  POST /backtest  → 历史回测",
             "  POST /explain   → SHAP 可解释性",
             "  POST /chat/stream → SSE 流式对话",
             "  GET  /           → 静态聊天 UI",
             "",
             "启动:",
             "  uvicorn ellectric.api.server:app",
             "  --host 0.0.0.0 --port 8000"],
         title_color=ACCENT)

# CLI
add_card(slide, Inches(4.6), Inches(1.5), Inches(4.0), Inches(4.0),
         "💻 Typer CLI 工具",
         [
             "5 个子命令:",
             "  forecast load 24",
             "    → 未来 24h 负荷预测",
             "  simulate summer_peak --days 7",
             "    → 夏季高峰 7 天仿真",
             "  backtest 2022-08-01 2022-08-31 ppo",
             "    → 8月 PPO 策略回测",
             "  explain xgboost 0",
             "    → 第 0 个样本 SHAP 解释",
             "  ask \"昨天峰值负荷？\"",
             "    → LLM 自然语言查询",
             "",
             "支持 --json 输出 + rich 表格"],
         title_color=ACCENT2)

# LLM
add_card(slide, Inches(8.9), Inches(1.5), Inches(3.9), Inches(4.0),
         "🤖 LangChain Agent",
         [
             "LLM: DeepSeek Chat API",
             "  兼容 OpenAI SDK 格式",
             "  model: deepseek-v4-flash",
             "  temperature: 0.3",
             "",
             "3 个 Tool 函数:",
             "  query_forecast()  ",
             "    → 调用 /predict API",
             "  run_simulation()  ",
             "    → 调用 /simulate API",
             "  run_backtest()    ",
             "    → 调用 /backtest API",
             "",
             "SSE 流式:",
             "  astream_events() → token 帧"],
         title_color=ORANGE)

add_card(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2),
         "🏗️ 三明治架构: API / CLI / LLM → handlers.py (4个函数) → Pipeline (12个模块) → 模型/数据文件",
         ["所有接入层通过统一的 Service 层 (handlers) 调用 Pipeline — 延迟导入安全，允许部分模块不可用"],
         title_color=ACCENT)

add_page_number(slide, 14)

# ═══════════════════════ SLIDE 15: 技术栈全景 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "技术栈全景图", "Ellectric 使用的所有技术及版本")

stacks = [
    ("📦 数据处理", ACCENT, ["pandas 3.0.x", "numpy ≥2.0.0", "pyarrow 22.0", "urllib+csv (stdlib)"]),
    ("🧠 机器学习", GREEN, ["scikit-learn 1.8.0", "xgboost 3.2.0", "shap", "joblib"]),
    ("🎮 强化学习", ACCENT3, ["stable-baselines3", "gymnasium", "PPO / SAC / TD3"]),
    ("📊 可视化", ORANGE, ["plotly 6.7.0", "Jupyter 1.1.1", "10 个渐进式 notebooks"]),
    ("🏭 仿真", ACCENT2, ["ASSUME 0.6.0", "Merit Order 内置引擎", "YAML 场景配置"]),
    ("🌐 Web + CLI", ACCENT, ["FastAPI 0.136.x", "Pydantic v2 (Rust后端)", "Typer", "uvicorn"]),
    ("🤖 LLM", GREEN, ["LangChain 1.3.x", "DeepSeek v4-flash", "httpx", "SSE 流式"]),
    ("📐 统计", ACCENT3, ["epftoolbox (git)", "Diebold-Mariano 检验", "Giacomini-White 检验"]),
]

for i, (name, color, items) in enumerate(stacks):
    row = i // 4
    col = i % 4
    left = Inches(0.5 + col * 3.15)
    top = Inches(1.5 + row * 2.7)
    item_text = "\n".join(f"  {item}" for item in items)
    add_card(slide, left, top, Inches(2.95), Inches(2.3),
             name, [item_text], title_color=color)

add_page_number(slide, 15)

# ═══════════════════════ SLIDE 16: 算法选择原因 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide)
add_title_bar(slide, "算法选择原因 — 为什么是这个组合？", "每个技术选型背后的设计决策")

decisions = [
    ("XGBoost 负荷预测", GREEN,
     "梯度提升树是表格数据的 SOTA。Kaggle 能源预测竞赛长期霸榜。原生处理缺失值，内置特征重要性。"
     "比深度学习更稳健 (小样本友好)，CPU 优化不需要 GPU。TimeSeriesSplit 防泄漏设计成熟。"),
    ("LEAR/Lasso 电价预测", ACCENT2,
     "LEAR 是电力日前价格预测的标准基准方法 (Lago et al., 2021)。Lasso 的 L1 正则化自动特征选择，"
     "系数可直接解释为边际经济影响 — 这在电价预测中比黑盒精度更重要。与 XGBoost 形成互补教学对比。"),
    ("PPO 强化学习 (首选)", ACCENT3,
     "PPO 是连续控制任务的工业标准 (OpenAI Five, ChatGPT RLHF)。On-policy → 稳定性好，不易发散。"
     "clipping 机制防止策略突变。SAC/TD3 作为备选算法提供对比学习价值。"),
    ("Plotly 而非 Matplotlib", ORANGE,
     "Plotly 生成交互式图表: 悬停查看精确数值、框选放大、双击重置 — 对电力数据探索至关重要。"
     "Jupyter 中自动渲染。时间序列叠加图 + 误差直方图 + P&L 曲线子图组合远超静态图片的效果。"),
    ("FastAPI + Pydantic v2", ACCENT,
     "FastAPI 是 Python 最快 REST 框架之一 (Starlette + Pydantic)。Pydantic v2 用 Rust 重写核心，"
     "数据校验零开销。自动生成 OpenAPI/Swagger 文档。异步支持 SSE 流式对话。"),
    ("TimeSeriesSplit 不随机", RED,
     "普通 KFold 随机分割会引入未来信息泄露 (用 2023 数据预测 2021)。TimeSeriesSplit 保证时序: "
     "永远过去预测未来。gap=24 参数额外保护 lag 特征不跨越训练/测试边界。"),
    ("三明治架构", ACCENT2,
     "API/CLI/LLM 三层平行接入同一 Service 层 → 避免代码重复。Service 层延迟导入 Pipeline → "
     "模块间无循环依赖，部分模块不可用时不影响其他功能。清晰的职责分离。"),
    ("异常值保留不删除", RED,
     "电力负荷尖峰是有效信号: 极端天气/寒潮/重大活动导致的负荷突变不应自动剔除。"
     "IQR 检测仅报告 → 数据质量透明，让使用者(人/模型)自行决定如何处理。"),
]

for i, (title, color, reason) in enumerate(decisions):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.5 + row * 1.35)
    add_card(slide, left, top, Inches(6.1), Inches(1.2),
             f"💡 {title}", [reason], title_color=color)

add_page_number(slide, 16)

# ═══════════════════════ SLIDE 17: 结束页 ═══════════════════════
slide = prs.slides.add_slide(BLANK)
add_bg(slide, BG_DARK)

add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
         "谢谢！", font_size=64, bold=True, color=WHITE)

add_text(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
         "Ellectric — AI + 电力交易技术学习平台", font_size=24, color=ACCENT)

add_text(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.8),
         "数据 → 特征 → 预测 → 仿真 → 交易 → 解释 → 问答  |  端到端技术闭环",
         font_size=14, color=MID_GREY)

# 核心要点回顾
add_text(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.4),
         "核心要点回顾", font_size=18, bold=True, color=ACCENT)

add_bullet_list(slide, Inches(1.5), Inches(5.3), Inches(10.3),
    [
        "电力数据: OWID 全球年数据 (自动拉取) + 中国现货电价小时数据 (本地 xlsx)",
        "负荷预测: XGBoost + TimeSeriesSplit 防时序泄露 + 持续法基线对比",
        "电价预测: LEAR (Lasso L1 正则化) — 可解释且精度经同行评审验证",
        "市场仿真: ASSUME 多智能体框架 + Merit Order 内置回退引擎",
        "交易策略: PPO/TD3/SAC 三种 RL 算法 + 多策略回测对比",
        "模型解释: SHAP TreeExplainer + DM/GW 统计显著性检验",
        "服务接口: FastAPI REST + Typer CLI + LangChain/DeepSeek 自然语言对话",
        "渐进式学习: 10 个 Jupyter notebooks 按 pipeline 顺序逐步推进",
    ],
    font_size=11,
)

add_text(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.3),
         "Python 3.11+ | XGBoost 3.2.0 | scikit-learn 1.8.0 | stable-baselines3 | FastAPI 0.136.x | LangChain 1.3.x",
         font_size=10, color=MID_GREY)

# ═══════════════════════ 保存 ═══════════════════════
output_dir = os.getcwd()
output_path = os.path.join(output_dir, "Ellectric_技术讲解.pptx")
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
